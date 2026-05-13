"""
core/document_reader.py
========================
Extracts plain text from PDF, TXT, DOCX, and PPTX files.
Returns a clean string ready for analyse_transcript().
"""

import io
from pathlib import Path


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract text from a document file.

    Args:
        file_bytes: Raw file content.
        filename:   Original filename (used to detect format).

    Returns:
        Extracted text as a single string.

    Raises:
        ValueError:  Unsupported file format.
        RuntimeError: Extraction failure.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".txt":
        return _extract_txt(file_bytes)
    elif ext == ".pdf":
        return _extract_pdf(file_bytes)
    elif ext == ".docx":
        return _extract_docx(file_bytes)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported format '{ext}'. Use: .txt, .pdf, .docx, .pptx")


def _extract_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore").strip()
    except Exception as exc:
        raise RuntimeError(f"TXT extraction failed: {exc}") from exc


def _extract_pdf(file_bytes: bytes) -> str:
    try:
        import fitz  # pymupdf
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page in doc:
            pages.append(page.get_text())
        return "\n".join(pages).strip()
    except Exception as exc:
        raise RuntimeError(f"PDF extraction failed: {exc}") from exc


def _extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs).strip()
    except Exception as exc:
        raise RuntimeError(f"DOCX extraction failed: {exc}") from exc


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text).strip()
    except Exception as exc:
        raise RuntimeError(f"PPTX extraction failed: {exc}") from exc