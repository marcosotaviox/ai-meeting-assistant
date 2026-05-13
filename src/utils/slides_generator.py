"""
utils/slides_generator.py
==========================
Renders a SlidesPlan into a PowerPoint file.
Layout decisions are made by slide_planner.py (Claude).
This module only handles rendering.
"""

import os
import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.core.analyser import MeetingInsights
from src.core.slide_planner import SlidesPlan, SlideContent, plan_slides
from src.utils.youtube_search import search_youtube

# ── Colour themes ─────────────────────────────────────────────────────────────
THEMES = {
    "cyan":   {"header": RGBColor(0x00, 0x56, 0x8A), "accent": RGBColor(0x00, 0xB8, 0xD9)},
    "purple": {"header": RGBColor(0x5B, 0x21, 0xB6), "accent": RGBColor(0xA7, 0x8B, 0xFA)},
    "green":  {"header": RGBColor(0x06, 0x5F, 0x46), "accent": RGBColor(0x34, 0xD3, 0x99)},
    "amber":  {"header": RGBColor(0x92, 0x40, 0x00), "accent": RGBColor(0xF5, 0x9E, 0x0B)},
}

NAVY      = RGBColor(0x0A, 0x0F, 0x1E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xCB, 0xD5, 0xE1)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
SUBTLE    = RGBColor(0x64, 0x74, 0x8B)
CARD_BG   = RGBColor(0xF1, 0xF5, 0xF9)
CARD_BDR  = RGBColor(0xE2, 0xE8, 0xF0)

# ── Dimensions ────────────────────────────────────────────────────────────────
W        = Inches(10)
H        = Inches(7.5)
HEADER_H = Inches(1.1)
BODY_TOP = Inches(1.1)
BODY_H   = Inches(6.0)
IMG_W    = Inches(4.0)
CONTENT_L= Inches(4.3)
CONTENT_W= Inches(5.5)
MARGIN   = Inches(0.35)


# ── Fetch helpers ─────────────────────────────────────────────────────────────

def _fetch_bytes(url: str) -> bytes | None:
    try:
        return requests.get(url, timeout=15).content
    except Exception as exc:
        print(f"Download failed: {exc}")
        return None


def _fetch_unsplash(query: str) -> bytes | None:
    key = os.environ.get("UNSPLASH_ACCESS_KEY")
    if not key or not query:
        return None
    try:
        r = requests.get(
            "https://api.unsplash.com/search/photos",
            params={"query": query, "per_page": 1, "orientation": "portrait"},
            headers={"Authorization": f"Client-ID {key}"},
            timeout=10,
        )
        url = r.json()["results"][0]["urls"]["regular"]
        return _fetch_bytes(url)
    except Exception as exc:
        print(f"Unsplash failed: {exc}")
        return None


# ── Drawing helpers ───────────────────────────────────────────────────────────

def _set_background(slide, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header(slide, title: str, header_color: RGBColor) -> None:
    rect = slide.shapes.add_shape(1, Inches(0), Inches(0), W, HEADER_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = header_color
    rect.line.fill.background()

    txBox = slide.shapes.add_textbox(MARGIN, Inches(0.15), W - MARGIN * 2, HEADER_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_text_box(slide, text, left, top, width, height,
                  size=14, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullets(slide, items: list[str], left, top, width, height,
                 size=14, color=DARK_TEXT) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_hyperlink(slide, text, url, left, top, width, height, size=11,
                   accent_color=None) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = accent_color or RGBColor(0x00, 0x56, 0x8A)
    run.font.underline = True

    from pptx.oxml.ns import qn
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    rId = slide.part.relate_to(
        url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True
    )
    hlinkClick.set(qn('r:id'), rId)


def _add_card(slide, title: str, body: str,
              left, top, width, height,
              accent_color: RGBColor = None) -> None:
    accent = accent_color or RGBColor(0x00, 0xB8, 0xD9)

    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = CARD_BDR

    bar = slide.shapes.add_shape(1, left, top, Inches(0.07), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    if title:
        _add_text_box(slide, title,
                      left + Inches(0.18), top + Inches(0.1),
                      width - Inches(0.25), Inches(0.35),
                      size=11, bold=True, color=accent)

    body_top = top + (Inches(0.38) if title else Inches(0.12))
    body_h   = height - (Inches(0.38) if title else Inches(0.12)) - Inches(0.1)
    _add_text_box(slide, body,
                  left + Inches(0.18), body_top,
                  width - Inches(0.25), body_h,
                  size=12, color=DARK_TEXT)


def _build_cards_grid(slide, items: list[dict],
                      accent_color: RGBColor = None) -> None:
    """Render cards in a smart grid filling the full slide body."""
    count = len(items)
    if count == 0:
        return

    pad    = MARGIN
    card_w = W - pad * 2
    total_h= BODY_H - pad * 2

    if count == 1:
        _add_card(slide, items[0].get("title",""), items[0].get("body",""),
                  pad, BODY_TOP + pad, card_w, total_h, accent_color)

    elif count == 2:
        card_h = (total_h - pad) / 2
        for i, item in enumerate(items):
            top = BODY_TOP + pad + i * (card_h + pad)
            _add_card(slide, item.get("title",""), item.get("body",""),
                      pad, top, card_w, card_h, accent_color)

    elif count == 3:
        card_h = (total_h - pad * 2) / 3
        for i, item in enumerate(items):
            top = BODY_TOP + pad + i * (card_h + pad)
            _add_card(slide, item.get("title",""), item.get("body",""),
                      pad, top, card_w, card_h, accent_color)

    else:
        col_w  = (card_w - pad) / 2
        rows   = (count + 1) // 2
        card_h = (total_h - pad * (rows - 1)) / rows
        for i, item in enumerate(items):
            col  = i % 2
            row  = i // 2
            left = pad + col * (col_w + pad)
            top  = BODY_TOP + pad + row * (card_h + pad)
            _add_card(slide, item.get("title",""), item.get("body",""),
                      left, top, col_w, card_h, accent_color)


# ── Layout renderers ──────────────────────────────────────────────────────────

def _render_cover(prs, slide_data: SlideContent, insights: MeetingInsights,
                  theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, NAVY)

    img_bytes = _fetch_unsplash(insights.title)
    if img_bytes:
        pic = slide.shapes.add_picture(
            io.BytesIO(img_bytes), Inches(0), Inches(0), W, H
        )
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        rect = slide.shapes.add_shape(1, Inches(0), Inches(0), W, H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.fill.background()

    type_colors = {
        "Meeting": "00E5FF", "Class": "A78BFA",
        "Webinar": "00B8D9", "Presentation": "F59E0B"
    }
    hex_c = type_colors.get(insights.content_type, "00E5FF")
    badge = RGBColor(int(hex_c[:2], 16), int(hex_c[2:4], 16), int(hex_c[4:], 16))

    _add_text_box(slide, f"[ {insights.content_type.upper()} ]",
                  Inches(0.5), Inches(1.8), W - Inches(1), Inches(0.6),
                  12, bold=True, color=badge, align=PP_ALIGN.CENTER)

    _add_text_box(slide, slide_data.title,
                  Inches(0.5), Inches(2.5), W - Inches(1), Inches(2.2),
                  38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, f"Tone: {insights.sentiment}",
                  Inches(0.5), Inches(4.9), W - Inches(1), Inches(0.5),
                  14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    _add_text_box(slide,
                  "Generated by MeetingMind AI  ·  Whisper + Claude",
                  Inches(0.5), Inches(6.9), W - Inches(1), Inches(0.4),
                  9, color=RGBColor(0x00, 0x90, 0xA8), align=PP_ALIGN.CENTER)


def _render_summary(prs, slide_data: SlideContent, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    if slide_data.bullets:
        _add_bullets(slide, slide_data.bullets,
                     MARGIN, BODY_TOP + MARGIN, W - MARGIN * 2, BODY_H - MARGIN,
                     size=15)
    elif slide_data.cards:
        _build_cards_grid(slide, slide_data.cards, theme["accent"])


def _render_bullets(prs, slide_data: SlideContent, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    if slide_data.image_query:
        img_bytes = _fetch_unsplash(slide_data.image_query)
        if img_bytes:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(0), BODY_TOP, IMG_W, BODY_H
            )
            _add_bullets(slide, slide_data.bullets,
                         CONTENT_L, BODY_TOP + MARGIN, CONTENT_W, BODY_H - MARGIN,
                         size=14)
            return

    _add_bullets(slide, slide_data.bullets,
                 MARGIN, BODY_TOP + MARGIN, W - MARGIN * 2, BODY_H - MARGIN,
                 size=15)


def _render_image_left(prs, slide_data: SlideContent, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    img_bytes = _fetch_unsplash(slide_data.image_query) if slide_data.image_query else None
    if img_bytes:
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            Inches(0), BODY_TOP, IMG_W, BODY_H
        )
        _add_bullets(slide, slide_data.bullets,
                     CONTENT_L, BODY_TOP + MARGIN, CONTENT_W, BODY_H - MARGIN,
                     size=14)
    else:
        _add_bullets(slide, slide_data.bullets,
                     MARGIN, BODY_TOP + MARGIN, W - MARGIN * 2, BODY_H - MARGIN,
                     size=15)


def _render_cards(prs, slide_data: SlideContent, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    if slide_data.cards:
        _build_cards_grid(slide, slide_data.cards, theme["accent"])
    elif slide_data.bullets:
        items = [{"title": "", "body": b} for b in slide_data.bullets]
        _build_cards_grid(slide, items, theme["accent"])


def _render_action_items(prs, slide_data: SlideContent,
                          insights: MeetingInsights, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    top = BODY_TOP + MARGIN
    for item in insights.action_items[:6]:
        row_h = Inches(0.75)

        rect = slide.shapes.add_shape(1, MARGIN, top, W - MARGIN * 2, row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = CARD_BG
        rect.line.color.rgb = CARD_BDR

        badge = slide.shapes.add_shape(1, MARGIN, top, Inches(1.8), row_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = theme["header"]
        badge.line.fill.background()

        _add_text_box(slide, item.owner,
                      MARGIN + Inches(0.05), top + Inches(0.15),
                      Inches(1.7), row_h,
                      10, bold=True, color=WHITE)

        _add_text_box(slide, item.task,
                      MARGIN + Inches(1.9), top + Inches(0.15),
                      Inches(6.0), row_h,
                      12, color=DARK_TEXT)

        _add_text_box(slide, item.deadline,
                      W - Inches(1.6), top + Inches(0.15),
                      Inches(1.4), row_h,
                      10, color=SUBTLE, align=PP_ALIGN.RIGHT)

        top += row_h + Inches(0.12)


def _render_youtube(prs, slide_data: SlideContent,
                    insights: MeetingInsights, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_header(slide, slide_data.title, theme["header"])

    queries = []
    if insights.key_points:
        queries.append(insights.key_points[0][:80])
    if insights.key_concepts:
        queries.append(insights.key_concepts[0].concept)

    videos = []
    for query in queries:
        results = search_youtube(query, max_results=1)
        videos.extend(results)
        if len(videos) >= 2:
            break

    if not videos:
        _add_text_box(slide, "No YouTube videos found for this content.",
                      MARGIN, BODY_TOP + MARGIN, W - MARGIN * 2, Inches(1),
                      14, color=SUBTLE)
        return

    positions = [(MARGIN, BODY_TOP + MARGIN), (Inches(5.15), BODY_TOP + MARGIN)]
    thumb_w   = Inches(4.6)
    thumb_h   = Inches(3.8)

    for idx, video in enumerate(videos[:2]):
        left, top = positions[idx]
        thumb_bytes = _fetch_bytes(video["thumbnail"])
        if thumb_bytes:
            slide.shapes.add_picture(
                io.BytesIO(thumb_bytes), left, top, thumb_w, thumb_h
            )
        _add_hyperlink(slide,
                       f"▶  {video['title'][:55]}",
                       video["url"],
                       left, top + thumb_h + Inches(0.1),
                       thumb_w, Inches(0.4),
                       size=11, accent_color=theme["accent"])


# ── Layout dispatcher ─────────────────────────────────────────────────────────

def _render_slide(prs, slide_data: SlideContent,
                  insights: MeetingInsights, theme: dict) -> None:
    """Dispatch to the correct layout renderer based on slide_data.layout."""
    layout = slide_data.layout.lower()

    if layout == "cover":
        _render_cover(prs, slide_data, insights, theme)
    elif layout == "summary":
        _render_summary(prs, slide_data, theme)
    elif layout == "bullets":
        _render_bullets(prs, slide_data, theme)
    elif layout == "image_left":
        _render_image_left(prs, slide_data, theme)
    elif layout == "cards":
        _render_cards(prs, slide_data, theme)
    elif layout == "action_items":
        _render_action_items(prs, slide_data, insights, theme)
    elif layout == "youtube":
        _render_youtube(prs, slide_data, insights, theme)
    else:
        # Fallback — treat as bullets
        _render_bullets(prs, slide_data, theme)


# ── Main generator ────────────────────────────────────────────────────────────

def generate_slides(insights: MeetingInsights, include_youtube: bool = False) -> bytes:
    """
    Generate a PowerPoint from MeetingInsights.

    Flow:
    1. plan_slides() — Claude plans content, layouts, speaker notes
    2. Render each slide using the planned layout

    Args:
        insights:         Structured insights from analyse_transcript().
        include_youtube:  If True, Claude may include a YouTube layout slide.

    Returns:
        PPTX file as bytes.
    """
    # Step 1 — Plan
    plan = plan_slides(insights, include_youtube=include_youtube)

    # Step 2 — Theme
    theme = THEMES.get(plan.color_theme, THEMES["cyan"])

    # Step 3 — Render
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for slide_data in plan.slides:
        # Skip YouTube slides if not requested
        if slide_data.layout == "youtube" and not include_youtube:
            continue
        _render_slide(prs, slide_data, insights, theme)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()